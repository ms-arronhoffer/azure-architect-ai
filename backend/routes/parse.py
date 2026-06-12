"""Extract text from uploaded DOCX, PDF, or PPTX files."""

import io
import zipfile

from fastapi import APIRouter, HTTPException, Request
from fastapi.responses import JSONResponse

router = APIRouter()


@router.post("/parse")
async def parse_document(request: Request) -> JSONResponse:
    filename = request.headers.get("X-Filename", "upload")
    data = await request.body()

    if not data:
        raise HTTPException(status_code=422, detail="Uploaded file is empty.")

    lower = filename.lower()
    if lower.endswith(".pdf"):
        text = _extract_pdf(data)
    elif lower.endswith(".docx"):
        text = _extract_docx(data, filename)
    elif lower.endswith(".pptx"):
        text = _extract_pptx(data)
    elif lower.endswith(".xlsx") or lower.endswith(".xls"):
        text = _extract_xlsx(data)
    else:
        raise HTTPException(status_code=415, detail="Only .pdf, .docx, .pptx, and .xlsx/.xls files are supported.")

    return JSONResponse({"text": text, "filename": filename})


def _extract_pdf(data: bytes) -> str:
    try:
        from pypdf import PdfReader
    except ImportError:
        raise HTTPException(status_code=500, detail="pypdf not installed — run: python -m pip install pypdf") from None
    try:
        reader = PdfReader(io.BytesIO(data))
        parts = [page.extract_text() or "" for page in reader.pages]
        return "\n\n".join(p for p in parts if p.strip())
    except Exception as exc:
        raise HTTPException(status_code=422, detail=f"Could not read PDF: {exc}") from exc


def _extract_docx(data: bytes, filename: str) -> str:
    if not zipfile.is_zipfile(io.BytesIO(data)):
        raise HTTPException(
            status_code=422,
            detail=(
                f"'{filename}' does not appear to be a valid .docx file "
                f"({len(data)} bytes received). "
                "If this is a legacy .doc file, open in Word and save as .docx."
            ),
        )
    try:
        from docx import Document
    except ImportError:
        raise HTTPException(status_code=500, detail="python-docx not installed.") from None
    try:
        doc = Document(io.BytesIO(data))
        parts = [para.text for para in doc.paragraphs if para.text.strip()]
        return "\n\n".join(parts)
    except Exception as exc:
        raise HTTPException(status_code=422, detail=f"Could not read DOCX: {exc}") from exc


def _extract_pptx(data: bytes) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        raise HTTPException(status_code=500, detail="python-pptx not installed.") from None
    try:
        prs = Presentation(io.BytesIO(data))
        parts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = [
                shape.text.strip()
                for shape in slide.shapes
                if hasattr(shape, "text") and shape.text.strip()
            ]
            if slide_texts:
                parts.append(f"[Slide {slide_num}]\n" + "\n".join(slide_texts))
        return "\n\n".join(parts)
    except Exception as exc:
        raise HTTPException(status_code=422, detail=f"Could not read PPTX: {exc}") from exc


def _extract_xlsx(data: bytes) -> str:
    # XLS magic bytes: D0 CF 11 E0 (old BIFF binary format)
    if data[:4] == b"\xd0\xcf\x11\xe0":
        return _extract_xls(data)

    try:
        import openpyxl
    except ImportError:
        raise HTTPException(status_code=500, detail="openpyxl not installed — run: python -m pip install openpyxl") from None
    try:
        wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        parts = []
        for sheet in wb.worksheets:
            rows = []
            for row in sheet.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(c for c in cells):
                    rows.append("\t".join(cells))
            if rows:
                parts.append(f"[Sheet: {sheet.title}]\n" + "\n".join(rows))
        wb.close()
        return "\n\n".join(parts)
    except Exception as exc:
        raise HTTPException(status_code=422, detail=f"Could not read XLSX: {exc}") from exc


def _extract_xls(data: bytes) -> str:
    try:
        import xlrd
    except ImportError:
        raise HTTPException(
            status_code=422,
            detail="File is in old Excel 97-2003 (.xls) format. Open it in Excel and save as .xlsx, or contact your admin to enable .xls support.",
        ) from None
    try:
        wb = xlrd.open_workbook(file_contents=data)
        parts = []
        for sheet in wb.sheets():
            rows = []
            for rx in range(sheet.nrows):
                cells = [str(sheet.cell_value(rx, cx)) for cx in range(sheet.ncols)]
                if any(c for c in cells):
                    rows.append("\t".join(cells))
            if rows:
                parts.append(f"[Sheet: {sheet.name}]\n" + "\n".join(rows))
        return "\n\n".join(parts)
    except Exception as exc:
        raise HTTPException(status_code=422, detail=f"Could not read XLS: {exc}") from exc
