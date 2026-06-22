"""Presentation outline (SSE) and PPTX build endpoints."""

import io
import json

from typing import Literal

from fastapi import APIRouter, File, Form, UploadFile
from fastapi.responses import Response, StreamingResponse
from pydantic import BaseModel

from services.openai_service import get_client, get_deployment
from services.pptx_service import build_presentation
from tools.tool_definitions import get_tools

router = APIRouter()


class OutlineRequest(BaseModel):
    topic: str
    audience: str = ""
    objectives: str = ""
    num_slides: int = 10
    conversation_context: str = ""


class BuildRequest(BaseModel):
    outline: dict
    theme: Literal["light", "dark"] = "dark"
    accent: str | None = None


def _sse(payload: dict) -> str:
    return f"data: {json.dumps(payload)}\n\n"


async def extract_text_from_files(files: list[UploadFile]) -> str:
    parts: list[str] = []
    for f in files:
        if not f.filename:
            continue
        data = await f.read()
        ext = f.filename.rsplit(".", 1)[-1].lower() if "." in f.filename else ""
        text = ""
        try:
            if ext == "pdf":
                from pypdf import PdfReader
                reader = PdfReader(io.BytesIO(data))
                text = "\n".join(p.extract_text() or "" for p in reader.pages)
            elif ext == "docx":
                from docx import Document  # type: ignore[import-untyped]
                doc = Document(io.BytesIO(data))
                text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
            elif ext == "xlsx":
                import openpyxl
                wb = openpyxl.load_workbook(io.BytesIO(data), data_only=True)
                lines: list[str] = []
                for ws in wb.worksheets:
                    lines.append(f"[Sheet: {ws.title}]")
                    for row in ws.iter_rows(values_only=True):
                        row_str = " | ".join(str(c) for c in row if c is not None)
                        if row_str.strip():
                            lines.append(row_str)
                text = "\n".join(lines)
            elif ext == "xls":
                import xlrd  # type: ignore[import-untyped]
                wb = xlrd.open_workbook(file_contents=data)
                lines = []
                for ws in wb.sheets():
                    lines.append(f"[Sheet: {ws.name}]")
                    for row_idx in range(ws.nrows):
                        row_str = " | ".join(
                            str(ws.cell_value(row_idx, col))
                            for col in range(ws.ncols)
                            if ws.cell_value(row_idx, col) != ""
                        )
                        if row_str.strip():
                            lines.append(row_str)
                text = "\n".join(lines)
            elif ext == "pptx":
                from pptx import Presentation  # type: ignore[import-untyped]
                prs = Presentation(io.BytesIO(data))
                slide_texts: list[str] = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_texts.append(shape.text)
                text = "\n".join(slide_texts)
            else:
                text = data.decode("utf-8", errors="replace")
        except Exception:
            text = data.decode("utf-8", errors="replace")

        if text.strip():
            parts.append(f"--- {f.filename} ---\n{text[:6000]}")

    return "\n\n".join(parts)


async def _generate(req: OutlineRequest, file_context: str = ""):
    client = get_client()
    deployment = get_deployment("architecture")
    outline_tool = get_tools("generate_deck_outline")[0]
    review_tool  = get_tools("review_deck_outline")[0]

    system = (
        "You are a senior Azure solutions architect and professional presentation designer. "
        "Create polished, customer-facing presentation decks with clear narrative arcs. "
        "Always vary slide layouts — mix content, two_column, quote_stat, and section_dividers."
    )
    user = (
        f"Create a {req.num_slides}-slide presentation about: {req.topic}\n"
        f"Target audience: {req.audience or 'technical and business stakeholders'}\n"
        f"Key objectives: {req.objectives or 'educate and inspire action'}\n\n"
        "Required structure: title slide → agenda → section dividers for major topics → "
        "body slides (varied layouts) → summary/key takeaways → references/next steps."
    )
    if req.conversation_context:
        user += (
            f"\n\nAdditional context from a coaching conversation — use this to inform "
            f"the deck's narrative, emphasis, and examples:\n\n{req.conversation_context[:3000]}"
        )
    if file_context:
        user += (
            f"\n\nSupporting documents provided by the user — incorporate relevant facts, "
            f"data, and context into the deck:\n\n{file_context[:6000]}"
        )

    # ── Call 1: generate outline ────────────────────────────────────────────
    resp1 = client.chat.completions.create(
        model=deployment,
        messages=[{"role": "system", "content": system}, {"role": "user", "content": user}],
        tools=[outline_tool],
        tool_choice={"type": "function", "function": {"name": "generate_deck_outline"}},
        max_completion_tokens=4000,
    )
    tc1 = resp1.choices[0].message.tool_calls[0]
    outline = json.loads(tc1.function.arguments)
    yield _sse({"type": "outline", "outline": outline})

    # ── Call 2: review + improve ────────────────────────────────────────────
    resp2 = client.chat.completions.create(
        model=deployment,
        messages=[
            {"role": "system", "content": system},
            {"role": "user", "content": user},
            resp1.choices[0].message,
            {"role": "tool", "tool_call_id": tc1.id, "content": json.dumps(outline)},
            {
                "role": "user",
                "content": (
                    "Critically review this outline. Evaluate: narrative flow, audience alignment, "
                    "content depth, layout variety, and storytelling impact. "
                    "Return recommendations AND a fully improved outline that addresses every issue."
                ),
            },
        ],
        tools=[review_tool],
        tool_choice={"type": "function", "function": {"name": "review_deck_outline"}},
        max_completion_tokens=5000,
    )
    tc2 = resp2.choices[0].message.tool_calls[0]
    review = json.loads(tc2.function.arguments)
    yield _sse({
        "type": "review",
        "overall_assessment": review.get("overall_assessment", ""),
        "recommendations":    review.get("recommendations", []),
        "improved_outline":   review.get("improved_outline", outline),
    })
    yield _sse({"type": "done"})


@router.post("/presentation/outline")
async def outline_endpoint(
    topic: str = Form(...),
    audience: str = Form(""),
    objectives: str = Form(""),
    num_slides: int = Form(10),
    conversation_context: str = Form(""),
    files: list[UploadFile] | None = File(default=None),
):
    req = OutlineRequest(
        topic=topic,
        audience=audience,
        objectives=objectives,
        num_slides=num_slides,
        conversation_context=conversation_context,
    )
    file_context = await extract_text_from_files(files) if files else ""
    return StreamingResponse(_generate(req, file_context=file_context), media_type="text/event-stream")


@router.post("/presentation/build")
async def build_endpoint(req: BuildRequest):
    pptx_bytes = build_presentation(req.outline, theme=req.theme, accent=req.accent)
    return Response(
        content=pptx_bytes,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": 'attachment; filename="presentation.pptx"'},
    )
