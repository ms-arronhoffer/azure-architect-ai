"""Theme-able PPTX generation service using python-pptx.

Supports a ``dark`` and ``light`` preset plus an optional accent override so the
customer can pick a world-class look that matches their brand. Every layout
renderer pulls colors from the active :class:`Theme` rather than module globals,
so adding a new preset is a one-line change.
"""

from dataclasses import dataclass, replace
from io import BytesIO

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def _hex(value: str) -> RGBColor:
    """Parse a ``#rrggbb`` (or ``rrggbb``) string into an RGBColor."""
    value = value.lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


@dataclass(frozen=True)
class Theme:
    """A self-contained color palette for a deck."""

    bg: RGBColor          # slide background
    card: RGBColor        # banded title bar / list rows
    accent: RGBColor      # primary brand accent (rules, numbers, dividers)
    accent_soft: RGBColor  # lighter accent for subtitles / stats
    title: RGBColor       # primary heading text
    body: RGBColor        # body / bullet text
    muted: RGBColor       # secondary / footer text
    divider: RGBColor     # thin separators
    on_accent: RGBColor   # text drawn on top of an accent fill
    divider_text: RGBColor  # subtitle text on a full-accent divider slide


DARK_THEME = Theme(
    bg=_hex("1A2028"),
    card=_hex("232B35"),
    accent=_hex("0078D4"),
    accent_soft=_hex("50C2FF"),
    title=_hex("FFFFFF"),
    body=_hex("E0E8F0"),
    muted=_hex("8098B0"),
    divider=_hex("2D3A47"),
    on_accent=_hex("FFFFFF"),
    divider_text=_hex("C0E0FF"),
)

LIGHT_THEME = Theme(
    bg=_hex("FFFFFF"),
    card=_hex("F2F5F9"),
    accent=_hex("0078D4"),
    accent_soft=_hex("106EBE"),
    title=_hex("1A2028"),
    body=_hex("2D3A47"),
    muted=_hex("5C6B7A"),
    divider=_hex("D6DEE8"),
    on_accent=_hex("FFFFFF"),
    divider_text=_hex("EAF3FF"),
)

_PRESETS = {"dark": DARK_THEME, "light": LIGHT_THEME}


def resolve_theme(name: str = "dark", accent: str | None = None) -> Theme:
    """Return a Theme preset, optionally overriding the accent color.

    ``accent`` may be a ``#rrggbb`` hex string; invalid values are ignored so a
    bad request never breaks the build.
    """
    theme = _PRESETS.get((name or "dark").lower(), DARK_THEME)
    if accent:
        try:
            theme = replace(theme, accent=_hex(accent))
        except (ValueError, IndexError):
            pass
    return theme


def _set_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _rect(slide, left, top, width, height, color: RGBColor):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _text(slide, text: str, left, top, width, height, size: int,
          bold=False, italic=False, color=None, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    if color is not None:
        run.font.color.rgb = color
    run.font.name = "Segoe UI"
    return tb


def _bullets(slide, items: list[str], left, top, width, height,
             size: int = 19, text_color=None):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(5)
        run = p.add_run()
        run.text = f"\u25aa  {item}"
        run.font.size = Pt(size)
        if text_color is not None:
            run.font.color.rgb = text_color
        run.font.name = "Segoe UI"
    return tb


def _notes(slide, text: str):
    slide.notes_slide.notes_text_frame.text = text


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _title_band(slide, t: "Theme", title: str):
    """Shared content-slide header: banded title bar, accent spine + rule."""
    _rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.0), t.card)
    _rect(slide, Inches(0), Inches(0), Inches(0.09), Inches(1.0), t.accent)
    _text(slide, title, Inches(0.35), Inches(0.15),
          Inches(12.5), Inches(0.7), size=28, bold=True, color=t.title)
    # Thin accent rule directly under the title band.
    _rect(slide, Inches(0), Inches(1.0), SLIDE_W, Pt(3), t.accent)


def _footer(slide, t: "Theme", deck_title: str, number: int):
    """Slide-number + deck-title footer for a consistent, polished look."""
    _text(slide, deck_title, Inches(0.4), Inches(7.06),
          Inches(10.5), Inches(0.34), size=10, color=t.muted)
    _text(slide, str(number), Inches(12.4), Inches(7.06),
          Inches(0.55), Inches(0.34), size=10, color=t.muted, align=PP_ALIGN.RIGHT)


# ── Layout renderers ─────────────────────────────────────────────────────────

def _title_slide(prs, t: "Theme", s: dict, deck_title: str, number: int):
    slide = _blank(prs)
    _set_bg(slide, t.bg)
    _rect(slide, Inches(0), Inches(0), Inches(0.09), SLIDE_H, t.accent)
    _rect(slide, Inches(0), Inches(6.85), SLIDE_W, Inches(0.65), t.accent)
    _text(slide, s.get("title", ""), Inches(0.5), Inches(1.6), Inches(9.5), Inches(2.2),
          size=48, bold=True, color=t.title)
    content = s.get("content") or []
    subtitle = content[0] if content else ""
    _text(slide, subtitle, Inches(0.5), Inches(3.95), Inches(9.5), Inches(1.6),
          size=24, color=t.accent_soft)
    for i in range(6):
        _rect(slide, Inches(10.8 + i * 0.42), Inches(2.2), Pt(7), Pt(7), t.accent)
    if s.get("speaker_notes"):
        _notes(slide, s["speaker_notes"])


def _agenda_slide(prs, t: "Theme", s: dict, deck_title: str, number: int):
    slide = _blank(prs)
    _set_bg(slide, t.bg)
    _rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.05), t.accent)
    _text(slide, s.get("title", "Agenda"), Inches(0.4), Inches(0.18),
          Inches(12), Inches(0.68), size=28, bold=True, color=t.on_accent)
    items = s.get("content") or []
    row_h = Inches(0.72)
    gap   = Inches(0.07)
    for idx, item in enumerate(items[:7]):
        top = Inches(1.22) + idx * (row_h + gap)
        _rect(slide, Inches(0.4), top, Inches(12.5), row_h, t.card)
        _rect(slide, Inches(0.4), top, Inches(0.54), row_h, t.accent)
        _text(slide, str(idx + 1), Inches(0.4), top + Inches(0.1),
              Inches(0.54), Inches(0.52), size=22, bold=True, color=t.on_accent,
              align=PP_ALIGN.CENTER)
        _text(slide, item, Inches(1.1), top + Inches(0.12),
              Inches(11.3), Inches(0.52), size=19, color=t.body)
    _footer(slide, t, deck_title, number)
    if s.get("speaker_notes"):
        _notes(slide, s["speaker_notes"])


def _divider_slide(prs, t: "Theme", s: dict, deck_title: str, number: int):
    slide = _blank(prs)
    _set_bg(slide, t.accent)
    _text(slide, s.get("title", ""), Inches(1.0), Inches(2.3), Inches(11.3), Inches(2.6),
          size=44, bold=True, color=t.on_accent, align=PP_ALIGN.CENTER)
    content = s.get("content") or []
    if content:
        _text(slide, content[0], Inches(1.5), Inches(5.1), Inches(10.3), Inches(1.3),
              size=20, color=t.divider_text, align=PP_ALIGN.CENTER)
    if s.get("speaker_notes"):
        _notes(slide, s["speaker_notes"])


def _content_slide(prs, t: "Theme", s: dict, deck_title: str, number: int):
    slide = _blank(prs)
    _set_bg(slide, t.bg)
    _title_band(slide, t, s.get("title", ""))
    _bullets(slide, s.get("content") or [], Inches(0.5), Inches(1.2),
             Inches(12.3), Inches(5.8), text_color=t.body)
    _footer(slide, t, deck_title, number)
    if s.get("speaker_notes"):
        _notes(slide, s["speaker_notes"])


def _two_column_slide(prs, t: "Theme", s: dict, deck_title: str, number: int):
    slide = _blank(prs)
    _set_bg(slide, t.bg)
    _title_band(slide, t, s.get("title", ""))
    _rect(slide, Inches(6.55), Inches(1.2), Inches(0.04), Inches(5.8), t.divider)
    _bullets(slide, s.get("content") or [], Inches(0.4), Inches(1.25),
             Inches(5.9), Inches(5.75), text_color=t.body)
    _bullets(slide, s.get("right_content") or [], Inches(6.8), Inches(1.25),
             Inches(6.1), Inches(5.75), text_color=t.body)
    _footer(slide, t, deck_title, number)
    if s.get("speaker_notes"):
        _notes(slide, s["speaker_notes"])


def _quote_stat_slide(prs, t: "Theme", s: dict, deck_title: str, number: int):
    slide = _blank(prs)
    _set_bg(slide, t.bg)
    content = s.get("content") or []
    _rect(slide, Inches(0.5), Inches(1.4), Inches(0.13), Inches(4.6), t.accent)
    main = content[0] if content else s.get("title", "")
    label = s.get("title", "")
    sub   = content[1] if len(content) > 1 else ""
    _text(slide, main, Inches(0.85), Inches(1.5), Inches(11.5), Inches(3.0),
          size=52, bold=True, color=t.accent_soft)
    _text(slide, label, Inches(0.85), Inches(4.8), Inches(11.0), Inches(1.0),
          size=22, color=t.title)
    if sub:
        _text(slide, sub, Inches(0.85), Inches(5.9), Inches(11.0), Inches(1.0),
              size=16, color=t.muted)
    _footer(slide, t, deck_title, number)
    if s.get("speaker_notes"):
        _notes(slide, s["speaker_notes"])


def _summary_slide(prs, t: "Theme", s: dict, deck_title: str, number: int):
    slide = _blank(prs)
    _set_bg(slide, t.bg)
    _title_band(slide, t, s.get("title", "Key Takeaways"))
    items  = s.get("content") or []
    row_h  = Inches(0.88)
    gap    = Inches(0.09)
    for idx, item in enumerate(items[:5]):
        top = Inches(1.22) + idx * (row_h + gap)
        _rect(slide, Inches(0.4), top, Inches(12.5), row_h, t.card)
        _rect(slide, Inches(0.4), top, Inches(0.09), row_h, t.accent)
        _text(slide, item, Inches(0.65), top + Inches(0.14),
              Inches(12.0), Inches(0.65), size=18, color=t.body)
    _footer(slide, t, deck_title, number)
    if s.get("speaker_notes"):
        _notes(slide, s["speaker_notes"])


def _references_slide(prs, t: "Theme", s: dict, deck_title: str, number: int):
    slide = _blank(prs)
    _set_bg(slide, t.bg)
    _rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.05), t.accent)
    _text(slide, s.get("title", "References & Resources"), Inches(0.4), Inches(0.18),
          Inches(12), Inches(0.68), size=28, bold=True, color=t.on_accent)
    _bullets(slide, s.get("content") or [], Inches(0.5), Inches(1.3),
             Inches(12.3), Inches(5.7), size=17, text_color=t.body)
    _footer(slide, t, deck_title, number)
    if s.get("speaker_notes"):
        _notes(slide, s["speaker_notes"])


_HANDLERS = {
    "title":           _title_slide,
    "agenda":          _agenda_slide,
    "section_divider": _divider_slide,
    "content":         _content_slide,
    "two_column":      _two_column_slide,
    "quote_stat":      _quote_stat_slide,
    "summary":         _summary_slide,
    "references":      _references_slide,
}


def build_presentation(outline: dict, theme: str = "dark",
                       accent: str | None = None) -> bytes:
    """Render an outline into a themed PPTX and return the raw bytes."""
    t = resolve_theme(theme, accent)
    deck_title = outline.get("deck_title", "")
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    for idx, slide_data in enumerate(outline.get("slides", []), start=1):
        handler = _HANDLERS.get(slide_data.get("layout", "content"), _content_slide)
        handler(prs, t, slide_data, deck_title, idx)
    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()
