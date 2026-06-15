"""
HLS CSA Org Tracker report generation service.

Accepts three input files (xlsx or csv):
  1. Manager List — TPID, Account Name, Azure CSA M
  2. ACR Data    — multi-month Power BI export (FY26-Jul … FY26-Jun)
  3. OU Data     — customer deployment inventory (us_hls.csv schema)

Automatically selects the last full calendar month's ACR column
(e.g. June 13 → FY26-May) and produces the 9-section markdown report.

All core risk logic is ported from model-iq/mcp/model-data/src/scorecard.ts.
"""
from __future__ import annotations

import contextlib
import csv
import io
import json
import re
from datetime import date, timedelta
from pathlib import Path
from typing import Any

_RETIREMENTS_FILE = Path(__file__).parent.parent / "data" / "model_iq" / "retirements.json"

RISK_WEIGHTS: dict[str, int] = {
    "overdue": 1_000_000,
    "critical": 500_000,
    "warning": 100_000,
    "watch": 10_000,
    "ok": 0,
}

RISK_EMOJI: dict[str, str] = {
    "overdue": "⚫",
    "critical": "🔴",
    "warning": "🟡",
    "watch": "🟢",
    "ok": "🟢",
}


# ─── File Loading ────────────────────────────────────────────────────────────


def _xlsx_to_csv_bytes(data: bytes) -> bytes:
    """Convert xlsx bytes → csv bytes (utf-8) so all values are uniform strings."""
    import csv as _csv

    import openpyxl  # type: ignore[import]

    wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    ws = wb.active
    out = io.StringIO()
    writer = _csv.writer(out)
    for row in ws.iter_rows(values_only=True):
        writer.writerow(["" if v is None else str(v) for v in row])
    wb.close()
    return out.getvalue().encode("utf-8")


def _to_raw_rows(data: bytes, filename: str) -> list[list[str]]:
    """Return raw rows as list[list[str]] regardless of source format."""
    if filename.lower().endswith(".xlsx"):
        data = _xlsx_to_csv_bytes(data)
    return _csv_to_raw_rows(data)


def _xlsx_to_raw_rows(data: bytes) -> list[list[Any]]:
    import openpyxl  # type: ignore[import]

    wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    ws = wb.active
    rows = [list(row) for row in ws.iter_rows(values_only=True)]
    wb.close()
    return rows


def _csv_to_raw_rows(data: bytes) -> list[list[str]]:
    try:
        text = data.decode("utf-8-sig")
    except UnicodeDecodeError:
        text = data.decode("cp1252")
    return [list(row) for row in csv.reader(io.StringIO(text))]


def load_file(data: bytes, filename: str) -> list[dict]:
    """Generic loader: xlsx or csv with row 1 as headers → list[dict]."""
    rows = _to_raw_rows(data, filename)
    if not rows:
        return []
    headers = [str(h).strip() if h is not None else "" for h in rows[0]]
    result = []
    for row in rows[1:]:
        d = {h: (str(row[i]).strip() if i < len(row) and row[i] is not None else "") for i, h in enumerate(headers)}
        result.append(d)
    return result


# ─── ACR Month Detection ─────────────────────────────────────────────────────


def get_last_full_month_column(today: date) -> str:
    """Return the Power BI column label for the last fully completed month.

    June 13 2026 → 'FY26-May'
    Microsoft FY: Jul 1 → Jun 30  (FY = calendar_year+1 when month ≥ 7)
    """
    last = today.replace(day=1) - timedelta(days=1)
    fy = last.year if last.month <= 6 else last.year + 1
    return f"FY{str(fy)[2:]}-{last.strftime('%b')}"


# ─── ACR Parsing ─────────────────────────────────────────────────────────────


def _parse_acr_value(raw: Any) -> float:
    if raw is None:
        return 0.0
    cleaned = str(raw).strip().replace("$", "").replace(",", "")
    if not cleaned or cleaned in ("-", "nan", "NaN"):
        return 0.0
    try:
        return float(cleaned)
    except ValueError:
        return 0.0


def parse_acr_data(data: bytes, filename: str, month_col: str) -> dict[str, float]:
    """Parse ACR multi-month file → {ACCOUNT_NAME_UPPER: monthly_acr_float}.

    ACR6.csv layout:
      Row 0: FiscalMonth,,FY26-Jul,…,FY26-May,…,Total
      Row 1: TPAccountName,ServiceCompGrouping,$ ACR,…   ← label row (skip)
      Row 2+: ABBOTT LABORATORIES,Total,"$37,027",…
    """
    rows = _to_raw_rows(data, filename)
    if not rows:
        return {}

    # Locate the header row (first row whose cell 0 is "FiscalMonth")
    header_row_idx: int | None = None
    for i, row in enumerate(rows):
        if row and str(row[0]).strip().lower() == "fiscalmonth":
            header_row_idx = i
            break

    if header_row_idx is None:
        raise ValueError("ACR file: cannot find a row with 'FiscalMonth' in column 0")

    headers = [str(c).strip() if c is not None else "" for c in rows[header_row_idx]]

    if month_col not in headers:
        fy_cols = [h for h in headers if h.startswith("FY")]
        raise ValueError(
            f"ACR column '{month_col}' not found in file. "
            f"Available FY columns: {fy_cols}"
        )

    col_idx = headers.index(month_col)

    result: dict[str, float] = {}
    for row in rows[header_row_idx + 2:]:  # skip header row + label row
        if not row or len(row) < 2:
            continue
        name = str(row[0]).strip() if row[0] is not None else ""
        grouping = str(row[1]).strip() if len(row) > 1 and row[1] is not None else ""
        if not name or grouping != "Total":
            continue
        val = row[col_idx] if len(row) > col_idx else ""
        result[name.upper()] = _parse_acr_value(val)

    return result


# ─── Org Mapping ─────────────────────────────────────────────────────────────


def _normalize_tpid(raw: Any) -> str:
    if raw is None or str(raw).strip() in ("", "NaN", "nan"):
        return ""
    try:
        n = float(str(raw))
        if n != n:  # NaN
            return ""
        return str(int(n))
    except (ValueError, TypeError):
        return str(raw).strip()


def _find_col(headers: list[str], candidates: list[str]) -> str | None:
    lower = {h.lower().strip(): h for h in headers}
    for c in candidates:
        if c.lower() in lower:
            return lower[c.lower()]
    return None


def build_org_map(rows: list[dict]) -> tuple[dict, dict[str, list[str]]]:
    """Build org_map and account_directors from Manager List rows.

    Returns:
      org_map            — {director_alias → {accounts: [{tpid, name}]}}
      account_directors  — {tpid → ordered list of director aliases}
    """
    if not rows:
        return {}, {}

    headers = list(rows[0].keys())
    tpid_col = _find_col(headers, ["TPID", "tpid", "Top Parent ID"])
    name_col = _find_col(headers, ["Account Name", "TP Name", "account name"])
    dir_col = _find_col(headers, ["Azure CSA M", "azure csa m", "Director", "Manager"])

    if not tpid_col or not name_col or not dir_col:
        raise ValueError(
            f"Manager List is missing required columns (TPID, Account Name, Azure CSA M). "
            f"Found: {headers}"
        )

    org_map: dict[str, dict] = {}
    account_directors: dict[str, list[str]] = {}
    seen_per_dir: dict[str, set[str]] = {}

    for row in rows:
        raw_tpid = row.get(tpid_col, "").strip()
        raw_name = row.get(name_col, "").strip()
        raw_dir = row.get(dir_col, "").strip()

        if not raw_tpid or not raw_dir:
            continue

        tpid = _normalize_tpid(raw_tpid)
        if not tpid:
            continue

        # Maintain insertion-order director list per account
        if tpid not in account_directors:
            account_directors[tpid] = []
        if raw_dir not in account_directors[tpid]:
            account_directors[tpid].append(raw_dir)

        # Director → accounts (dedup by TPID)
        if raw_dir not in org_map:
            org_map[raw_dir] = {"accounts": []}
            seen_per_dir[raw_dir] = set()

        if tpid not in seen_per_dir[raw_dir]:
            seen_per_dir[raw_dir].add(tpid)
            org_map[raw_dir]["accounts"].append({"tpid": tpid, "name": raw_name})

    return org_map, account_directors


# ─── Deployment Parsing ──────────────────────────────────────────────────────


def parse_deployments(rows: list[dict]) -> list[dict]:
    result = []
    for row in rows:
        dep = dict(row)
        dep["_tpid"] = _normalize_tpid(dep.get("TPID", ""))
        result.append(dep)
    return result


def build_tpid_index(deployments: list[dict]) -> dict[str, list[dict]]:
    index: dict[str, list[dict]] = {}
    for dep in deployments:
        tpid = dep.get("_tpid", "")
        if not tpid:
            continue
        index.setdefault(tpid, []).append(dep)
    return index


def build_name_to_tpid(deployments: list[dict]) -> dict[str, str]:
    mapping: dict[str, str] = {}
    for dep in deployments:
        name = dep.get("TP Name", "").strip().upper()
        tpid = dep.get("_tpid", "")
        if name and tpid and name not in mapping:
            mapping[name] = tpid
    return mapping


# ─── Retirement Lookup ───────────────────────────────────────────────────────


def load_retirements() -> tuple[dict[str, date], dict[str, str]]:
    """Load retirements.json → (lookup, replacements).

    lookup      — {model_id_lower: earliest_date, 'model|version': date}
    replacements — {model_id_lower: cleaned_replacement_hint}
    """
    if not _RETIREMENTS_FILE.exists():
        return {}, {}

    raw = json.loads(_RETIREMENTS_FILE.read_text(encoding="utf-8"))
    models = raw.get("models", raw) if isinstance(raw, dict) else raw

    lookup: dict[str, date] = {}
    replacements: dict[str, str] = {}

    for r in models:
        model_id = (r.get("modelId") or "").lower().strip()
        if not model_id:
            continue

        ret_str = (r.get("retirementDate") or "").strip()
        if not ret_str:
            continue
        cleaned_date = re.sub(r"^No earlier than\s+", "", ret_str, flags=re.IGNORECASE).strip()[:10]
        try:
            ret_date = date.fromisoformat(cleaned_date)
        except ValueError:
            continue

        # Version-specific key
        version = str(r.get("version") or "").strip()
        if version and version.lower() != "none" and version != "null":
            lookup[f"{model_id}|{version}"] = ret_date

        # Model-level key: keep earliest date (most conservative)
        existing = lookup.get(model_id)
        if existing is None or ret_date < existing:
            lookup[model_id] = ret_date

        # Replacement hint (first one wins)
        hints = r.get("replacementHints") or []
        if hints and model_id not in replacements:
            hint = re.sub(r"\s+version:.*$", "", str(hints[0]), flags=re.IGNORECASE).strip()
            if hint:
                replacements[model_id] = hint

    return lookup, replacements


# ─── Risk Classification (port of scorecard.ts) ──────────────────────────────


def parse_token_tier(val: Any) -> int:
    """-1 = unknown, 0 = none, 1 = low, 2 = medium, 3 = high, 4 = very-high."""
    if val is None:
        return -1
    s = str(val).strip()
    if not s or s == "NaN":
        return -1
    if s.startswith("0"):
        return 0
    if s.startswith("1"):
        return 1
    if s.startswith("2"):
        return 2
    if s.startswith("3"):
        return 3
    if s.startswith("4"):
        return 4
    return -1


def is_ptu(dep: dict) -> bool:
    return "provision" in str(dep.get("OfferingName", "")).lower()


def is_active_deployment(dep: dict) -> bool:
    """Active = Week-1 Medium+ (tier ≥ 2) OR PTU provisioned capacity."""
    return is_ptu(dep) or parse_token_tier(dep.get("Tokens Week-1")) >= 2


def _epoch_day(d: date) -> int:
    from datetime import datetime
    return (datetime(d.year, d.month, d.day) - datetime(1970, 1, 1)).days


def classify_risk(
    model: str,
    version: str,
    ret_lookup: dict[str, date],
    pbi_date: str,
    today: date,
) -> tuple[str, int | None]:
    """Classify a single deployment's retirement risk.

    Returns (level, days_remaining).
    Priority: version-specific lookup > model-level lookup > PBI date field.
    """
    key = (model or "").lower().strip()
    if not key:
        return "ok", None

    ret_date: date | None = None
    if version:
        ret_date = ret_lookup.get(f"{key}|{version.strip()}")
    if ret_date is None:
        ret_date = ret_lookup.get(key)
    if ret_date is None and pbi_date:
        cleaned = re.sub(r"^No earlier than\s+", "", pbi_date, flags=re.IGNORECASE).strip()[:10]
        with contextlib.suppress(ValueError):
            ret_date = date.fromisoformat(cleaned)

    if ret_date is None:
        return "ok", None

    days = _epoch_day(ret_date) - _epoch_day(today)

    if days <= 0:
        return "overdue", days
    if days <= 90:
        return "critical", days
    if days <= 180:
        return "warning", days
    if days <= 365:
        return "watch", days
    return "ok", days


def compute_account_risk(
    name: str,
    tpid: str,
    deps: list[dict],
    ret_lookup: dict[str, date],
    today: date,
    monthly_acr: float = 0.0,
) -> dict:
    """Aggregate risk across all deployments for one account."""
    overdue = critical = warning = watch = 0
    min_days: int | None = None
    risk_models: set[str] = set()
    active_at_risk = 0

    for dep in deps:
        model = str(dep.get("Model", "")).strip()
        version = str(dep.get("Version") or "").strip()
        pbi_date = str(dep.get("Retirement Date") or "").strip()
        if pbi_date == "NaN":
            pbi_date = ""

        level, days = classify_risk(model, version, ret_lookup, pbi_date, today)

        if level == "overdue":
            overdue += 1
            risk_models.add(model)
        elif level == "critical":
            critical += 1
            risk_models.add(model)
        elif level == "warning":
            warning += 1
            risk_models.add(model)
        elif level == "watch":
            watch += 1
            risk_models.add(model)

        if level in ("overdue", "critical", "warning") and is_active_deployment(dep):
            active_at_risk += 1

        if days is not None and level != "ok" and (min_days is None or days < min_days):
            min_days = days

    if overdue > 0:
        acct_level = "overdue"
    elif critical > 0:
        acct_level = "critical"
    elif warning > 0:
        acct_level = "warning"
    elif watch > 0:
        acct_level = "watch"
    else:
        acct_level = "ok"

    priority_score = monthly_acr + RISK_WEIGHTS[acct_level]

    return {
        "name": name,
        "tpid": tpid,
        "totalDeployments": len(deps),
        "overdue": overdue,
        "critical": critical,
        "warning": warning,
        "watch": watch,
        "level": acct_level,
        "minDays": min_days,
        "atRiskModels": sorted(risk_models),
        "monthlyAcr": monthly_acr,
        "annualAcr": monthly_acr * 12,
        "priorityScore": priority_score,
        "activeAtRisk": active_at_risk,
        "directors": [],  # populated by build_org_scorecard
    }


# ─── Org Scorecard ───────────────────────────────────────────────────────────


def build_org_scorecard(
    org_map: dict,
    account_directors: dict[str, list[str]],
    tpid_index: dict[str, list[dict]],
    ret_lookup: dict[str, date],
    acr_by_name: dict[str, float],
    name_to_tpid: dict[str, str],
    today: date,
) -> dict:
    """Build the full org scorecard matching generateOrgScorecard() from scorecard.ts."""
    # Match ACR account names → TPIDs
    acr_map: dict[str, float] = {}
    for name_upper, acr in acr_by_name.items():
        tpid = name_to_tpid.get(name_upper)
        if tpid:
            acr_map[tpid] = acr_map.get(tpid, 0.0) + acr

    # Process all accounts, dedup by TPID (keep highest priority score)
    all_accounts_map: dict[str, dict] = {}
    for director, data in org_map.items():
        for acct in data["accounts"]:
            tpid = acct["tpid"]
            deps = tpid_index.get(tpid, [])
            if not deps:
                continue
            monthly_acr = acr_map.get(tpid, 0.0)
            risk = compute_account_risk(acct["name"], tpid, deps, ret_lookup, today, monthly_acr)
            risk["directors"] = account_directors.get(tpid, [director])

            existing = all_accounts_map.get(tpid)
            if existing is None or risk["priorityScore"] > existing["priorityScore"]:
                all_accounts_map[tpid] = risk

    all_accounts = sorted(all_accounts_map.values(), key=lambda a: -a["priorityScore"])
    top5 = all_accounts[:5]

    # Director-group summary (Section 3): group accounts by combined director set
    dir_group: dict[str, list[dict]] = {}
    for acct in all_accounts:
        dir_key = ", ".join(acct["directors"])
        dir_group.setdefault(dir_key, []).append(acct)

    director_summaries = []
    for dir_key, accts in dir_group.items():
        od = sum(1 for a in accts if a["level"] == "overdue")
        cr = sum(1 for a in accts if a["level"] == "critical")
        wn = sum(1 for a in accts if a["level"] == "warning")
        dir_level = "overdue" if od > 0 else "critical" if cr > 0 else "warning" if wn > 0 else "ok"
        director_summaries.append({
            "director": dir_key,
            "accounts": len(accts),
            "deployments": sum(a["totalDeployments"] for a in accts),
            "monthlyAcr": sum(a["monthlyAcr"] for a in accts),
            "overdue": od,
            "critical": cr,
            "warning": wn,
            "level": dir_level,
        })
    director_summaries.sort(key=lambda d: -d["monthlyAcr"])

    # Totals
    total_monthly = sum(a["monthlyAcr"] for a in all_accounts)
    od_accounts = [a for a in all_accounts if a["level"] == "overdue"]
    cr_accounts = [a for a in all_accounts if a["level"] == "critical"]
    wn_accounts = [a for a in all_accounts if a["level"] == "warning"]

    # Count unique directors in org map vs those with deployed accounts
    all_dir_with_deps = {d for acct in all_accounts for d in acct["directors"]}

    return {
        "allAccounts": all_accounts,
        "top5": top5,
        "directorSummaries": director_summaries,
        "totals": {
            "directorsInOrg": len(org_map),
            "directorsWithDeployments": len(all_dir_with_deps),
            "accountsWithDeployments": len(all_accounts),
            "totalDeployments": sum(a["totalDeployments"] for a in all_accounts),
            "totalMonthlyAcr": total_monthly,
            "overdue": len(od_accounts),
            "critical": len(cr_accounts),
            "warning": len(wn_accounts),
            "ok": len(all_accounts) - len(od_accounts) - len(cr_accounts) - len(wn_accounts),
            "overdueAcr": sum(a["monthlyAcr"] for a in od_accounts),
            "criticalAcr": sum(a["monthlyAcr"] for a in cr_accounts),
            "warningAcr": sum(a["monthlyAcr"] for a in wn_accounts),
        },
    }


def build_model_fleet_summary(
    tpid_index: dict[str, list[dict]],
    ret_lookup: dict[str, date],
    ret_replacements: dict[str, str],
    today: date,
) -> list[dict]:
    """Fleet-wide model retirement exposure (≤180 days, Section 4)."""
    model_data: dict[str, dict] = {}

    for deps in tpid_index.values():
        for dep in deps:
            model = str(dep.get("Model", "")).strip()
            version = str(dep.get("Version") or "").strip()
            pbi_date = str(dep.get("Retirement Date") or "").strip()
            if pbi_date == "NaN":
                pbi_date = ""
            tpid = dep.get("_tpid", "")

            level, days = classify_risk(model, version, ret_lookup, pbi_date, today)
            if level not in ("overdue", "critical", "warning") or days is None:
                continue

            key = model.lower()
            if key not in model_data:
                # Model-level retirement date for display
                ret_date = ret_lookup.get(key)
                model_data[key] = {
                    "model": model,
                    "retirementDate": ret_date.isoformat() if ret_date else "",
                    "days": days,
                    "deploys": 0,
                    "accounts": set(),
                    "replacement": ret_replacements.get(key, ""),
                }

            # Track most-urgent days
            if days < model_data[key]["days"]:
                model_data[key]["days"] = days

            model_data[key]["deploys"] += 1
            model_data[key]["accounts"].add(tpid)

    result = []
    for stat in model_data.values():
        result.append({
            "model": stat["model"],
            "retirementDate": stat["retirementDate"],
            "days": stat["days"],
            "hlsDeploys": stat["deploys"],
            "hlsAccounts": len(stat["accounts"]),
            "replacement": stat["replacement"] or "—",
        })

    result.sort(key=lambda m: m["days"])
    return result


# ─── Report Rendering ────────────────────────────────────────────────────────


def _fmt_acr(amount: float) -> str:
    return f"${amount:,.0f}"


def _fmt_fy(monthly: float) -> str:
    annual = monthly * 12
    if annual >= 1_000_000:
        return f"${annual / 1_000_000:.1f}M"
    if annual >= 1_000:
        return f"${annual / 1_000:.0f}K"
    return f"${annual:.0f}"


def _fmt_days(days: int | None) -> str:
    if days is None:
        return "—"
    return f"{days}d"


def _key_models(models: list[str], n: int = 3) -> str:
    if not models:
        return "—"
    head = models[:n]
    rest = len(models) - n
    s = ", ".join(head)
    return f"{s} +{rest}" if rest > 0 else s


def _ordinal(n: int) -> str:
    if 11 <= n % 100 <= 13:
        return f"{n}th"
    return f"{n}{('th', 'st', 'nd', 'rd', 'th', 'th', 'th', 'th', 'th', 'th')[n % 10]}"


def render_report(
    org_scorecard: dict,
    model_summary: list[dict],
    today: date,
    acr_month_label: str,
    manager_list_name: str,
    acr_file_name: str,
    ou_file_name: str,
) -> str:
    """Render the 9-section markdown report matching hls-csa-org-tracker.md."""
    t = org_scorecard["totals"]
    all_accounts: list[dict] = org_scorecard["allAccounts"]
    top5: list[dict] = org_scorecard["top5"]
    dir_sums: list[dict] = org_scorecard["directorSummaries"]

    L: list[str] = []

    # ── Header ──────────────────────────────────────────────────────────────
    L += [
        "# HLS CSA Org Tracker — Model Retirement Risk Report",
        "",
        f"- **Generated:** {today.isoformat()}",
        "- **Scope:** US HLS Field Accountability Unit",
        f"- **Source:** CSV org mapping ({manager_list_name})",
        f"- **ACR Source:** Power BI actual metered consumption ({acr_month_label})",
        "",
    ]

    # ── Section 1: Executive Summary ────────────────────────────────────────
    L += [
        "## 1. Executive Summary",
        "",
        "| Metric | Value |",
        "|--------|------:|",
        f"| Directors in org mapping | {t['directorsInOrg']} |",
        f"| Directors with deployed accounts | {t['directorsWithDeployments']} |",
        f"| Accounts with deployments | {t['accountsWithDeployments']} |",
        f"| Total deployments | {t['totalDeployments']:,} |",
        f"| **Total AI ACR** | **{_fmt_acr(t['totalMonthlyAcr'])}/mo** ({_fmt_fy(t['totalMonthlyAcr'])}/yr) |",
        "",
        "| Risk Level | Accounts | AI ACR at Risk |",
        "|------------|---------:|---------------:|",
        f"| ⚫ OVERDUE | {t['overdue']} | {_fmt_acr(t['overdueAcr'])}/mo |",
        f"| 🔴 CRITICAL | {t['critical']} | {_fmt_acr(t['criticalAcr'])}/mo |",
        f"| 🟡 WARNING | {t['warning']} | {_fmt_acr(t['warningAcr'])}/mo |",
        f"| 🟢 OK | {t['ok']} | — |",
        "",
        f"> **⚠️ {_fmt_acr(t['totalMonthlyAcr'])}/mo "
        f"({_fmt_fy(t['totalMonthlyAcr'])}/yr) of AI ACR is at risk** "
        f"across {t['accountsWithDeployments']} accounts.",
        "",
    ]

    # ── Section 2: Top 5 Priority Accounts ──────────────────────────────────
    L += [
        "## 2. 🏆 Top 5 Priority Accounts (by AI ACR)",
        "",
        "| Rank | Account | Director | AI ACR/mo | FY Est. | Deployments | ⚫ OD | 🔴 CR | 🟡 WN | Min Days | Risk |",
        "|:----:|---------|----------|----------:|--------:|:-----------:|:-----:|:-----:|:-----:|:--------:|:----:|",
    ]
    for i, a in enumerate(top5, 1):
        L.append(
            f"| {i} | **{a['name']}** | {', '.join(a['directors'])} | "
            f"{_fmt_acr(a['monthlyAcr'])} | {_fmt_fy(a['monthlyAcr'])} | "
            f"{a['totalDeployments']:,} | {a['overdue']} | {a['critical']} | "
            f"{a['warning']} | {_fmt_days(a['minDays'])} | {RISK_EMOJI[a['level']]} |"
        )
    L.append("")

    for i, a in enumerate(all_accounts[5:10], 6):
        L.append(
            f"### {_ordinal(i)}: {a['name']} — "
            f"{_fmt_acr(a['monthlyAcr'])}/mo | {a['totalDeployments']} deployments | "
            f"{_fmt_days(a['minDays'])} {RISK_EMOJI[a['level']]}"
        )
    L.append("")

    # ── Section 3: Director Summary ──────────────────────────────────────────
    L += [
        "## 3. Director Summary",
        "",
        "| Director | Accounts | Deployments | AI ACR/mo | ⚫ OD | 🔴 CR | 🟡 WN | Risk |",
        "|----------|:--------:|------------:|----------:|:-----:|:-----:|:-----:|:----:|",
    ]
    for ds in dir_sums:
        L.append(
            f"| {ds['director']} | {ds['accounts']} | {ds['deployments']:,} | "
            f"{_fmt_acr(ds['monthlyAcr'])} | {ds['overdue']} | {ds['critical']} | "
            f"{ds['warning']} | {RISK_EMOJI[ds['level']]} |"
        )
    L.append("")

    # ── Section 4: Models Retiring ≤180 Days ────────────────────────────────
    L += [
        "## 4. Models Retiring ≤180 Days (HLS Exposure)",
        "",
        "| Model | Retires | Days | HLS Deploys | HLS Accounts | Replacement |",
        "|-------|---------|:----:|------------:|-------------:|-------------|",
    ]
    for m in model_summary:
        L.append(
            f"| {m['model']} | {m['retirementDate']} | {m['days']}d | "
            f"{m['hlsDeploys']:,} | {m['hlsAccounts']:,} | {m['replacement']} |"
        )
    L.append("")

    # ── Appendix helper ──────────────────────────────────────────────────────
    def _appendix_table(accounts: list[dict], risk_col: str, count_fn, emoji: str) -> list[str]:
        rows: list[str] = []
        for a in accounts:
            rows.append(
                f"| **{a['name']}** | {', '.join(a['directors'])} | "
                f"{_fmt_acr(a['monthlyAcr'])} | {a['totalDeployments']:,} | "
                f"{count_fn(a)} | {_fmt_days(a['minDays'])} | "
                f"{_key_models(a['atRiskModels'])} |"
            )
        return rows

    # ── Section 5: Appendix A — Overdue ─────────────────────────────────────
    od_accts = [a for a in all_accounts if a["level"] == "overdue"]
    L += [
        f"## 5. Appendix A — ⚫ Overdue Accounts ({len(od_accts)})",
        "",
        "| Account | Director | AI ACR/mo | Deployments | ⚫ Overdue | Min Days | Key Models |",
        "|---------|----------|----------:|------------:|-----------:|:--------:|------------|",
    ]
    L += _appendix_table(od_accts, "overdue", lambda a: a["overdue"], "⚫")
    L.append("")

    # ── Section 6: Appendix B — Critical ────────────────────────────────────
    cr_accts = [a for a in all_accounts if a["level"] == "critical"]
    L += [
        f"## 6. Appendix B — 🔴 Critical Accounts ({len(cr_accts)})",
        "",
        "| Account | Director | AI ACR/mo | Deployments | 🔴 Critical | Min Days | Key Models |",
        "|---------|----------|----------:|------------:|------------:|:--------:|------------|",
    ]
    L += _appendix_table(cr_accts, "critical", lambda a: a["critical"], "🔴")
    L.append("")

    # ── Section 7: Appendix C — Warning ─────────────────────────────────────
    wn_accts = [a for a in all_accounts if a["level"] == "warning"]
    L += [
        f"## 7. Appendix C — 🟡 Warning Accounts ({len(wn_accts)})",
        "",
        "| Account | Director | AI ACR/mo | Deployments | 🟡 Warning | Min Days | Key Models |",
        "|---------|----------|----------:|------------:|-----------:|:--------:|------------|",
    ]
    L += _appendix_table(wn_accts, "warning", lambda a: a["warning"], "🟡")
    L.append("")

    # ── Section 8: Appendix D — All Accounts ────────────────────────────────
    L += [
        "## 8. Appendix D — All Accounts (ACR-Ranked)",
        "",
        "| Account | Director | AI ACR/mo | FY Est. | Deployments | Risk | ⚫ OD | 🔴 CR | 🟡 WN | Min Days |",
        "|---------|----------|----------:|--------:|:-----------:|:----:|:-----:|:-----:|:-----:|:--------:|",
    ]
    for a in all_accounts:
        L.append(
            f"| {a['name']} | {', '.join(a['directors'])} | "
            f"{_fmt_acr(a['monthlyAcr'])} | {_fmt_fy(a['monthlyAcr'])} | "
            f"{a['totalDeployments']:,} | {RISK_EMOJI[a['level']]} | "
            f"{a['overdue']} | {a['critical']} | {a['warning']} | "
            f"{_fmt_days(a['minDays'])} |"
        )
    L.append("")

    # ── Section 9: Data Sources & Notes ─────────────────────────────────────
    L += [
        "## 9. Data Sources & Notes",
        "",
        "| Source | File | Freshness |",
        "|--------|------|-----------|",
        f"| Org Mapping | {manager_list_name} | CSV input |",
        f"| Retirements | retirements.json | {today.isoformat()} |",
        f"| Customer Deployments | {ou_file_name} | CSV input |",
        f"| **AI ACR** | **{acr_file_name}** | "
        f"**Power BI actual ({acr_month_label}) — refreshed {today.isoformat()}** |",
        "",
        "### Risk Thresholds",
        "",
        "| Level | Condition |",
        "|-------|-----------|",
        "| ⚫ OVERDUE | Days remaining ≤ 0 (past retirement date) |",
        "| 🔴 CRITICAL | Days remaining ≤ 90 |",
        "| 🟡 WARNING | Days remaining ≤ 180 |",
        "| 🟢 OK | Days remaining > 180 |",
        "",
        "### Notes",
        "",
        "- **ACR figures are Power BI actual metered consumption**. "
        "This is the source of truth for revenue.",
        "- Accounts are ranked by **AI ACR first**, then risk urgency as tiebreaker.",
        f"- ACR column used: **{acr_month_label}** "
        f"(last full month as of {today.isoformat()})",
    ]

    return "\n".join(L) + "\n"


# ─── Main Entry Point ────────────────────────────────────────────────────────


def compute_org_data(
    manager_list_data: bytes,
    manager_list_name: str,
    acr_data: bytes,
    acr_name: str,
    ou_data: bytes,
    ou_name: str,
    today: date,
) -> tuple[dict, list[dict], str]:
    """Load and process all input files. Returns (org_scorecard, model_summary, month_col)."""
    month_col = get_last_full_month_column(today)
    ret_lookup, ret_replacements = load_retirements()

    ml_rows = load_file(manager_list_data, manager_list_name)
    org_map, account_directors = build_org_map(ml_rows)

    acr_by_name = parse_acr_data(acr_data, acr_name, month_col)

    ou_rows = load_file(ou_data, ou_name)
    deployments = parse_deployments(ou_rows)
    tpid_index = build_tpid_index(deployments)
    name_to_tpid = build_name_to_tpid(deployments)

    org_scorecard = build_org_scorecard(
        org_map, account_directors, tpid_index,
        ret_lookup, acr_by_name, name_to_tpid, today,
    )
    model_summary = build_model_fleet_summary(tpid_index, ret_lookup, ret_replacements, today)
    return org_scorecard, model_summary, month_col


def generate_org_report(
    manager_list_data: bytes,
    manager_list_name: str,
    acr_data: bytes,
    acr_name: str,
    ou_data: bytes,
    ou_name: str,
    today: date | None = None,
) -> str:
    """Orchestrate all phases and return the 9-section markdown report."""
    if today is None:
        today = date.today()
    org_scorecard, model_summary, month_col = compute_org_data(
        manager_list_data, manager_list_name,
        acr_data, acr_name,
        ou_data, ou_name,
        today,
    )
    return render_report(
        org_scorecard, model_summary, today, month_col,
        manager_list_name, acr_name, ou_name,
    )


def generate_recommendations(
    org_scorecard: dict,
    model_summary: list[dict],
    today: date,
) -> str:
    """Generate AI-powered recommendations from the org scorecard (mirrors model-iq advisor)."""
    from services.openai_service import get_client, get_deployment  # type: ignore[import]

    totals = org_scorecard.get("totals", {})
    all_accounts = org_scorecard.get("allAccounts", [])
    director_summaries = org_scorecard.get("directorSummaries", [])

    urgent = [a for a in all_accounts if a.get("level") in ("overdue", "critical")][:25]
    warning_accts = [a for a in all_accounts if a.get("level") == "warning"][:15]

    def _a(val: float) -> str:
        return f"${val:,.0f}"

    ctx: list[str] = [
        f"**Date:** {today.isoformat()}",
        f"**Org:** {totals.get('accountsWithDeployments', 0)} accounts | "
        f"{totals.get('directorsWithDeployments', 0)} directors | "
        f"{_a(totals.get('totalMonthlyAcr', 0))}/mo ACR | "
        f"{totals.get('totalDeployments', 0)} deployments",
        f"**Risk:** ⚫ {totals.get('overdue', 0)} overdue | "
        f"🔴 {totals.get('critical', 0)} critical | "
        f"🟡 {totals.get('warning', 0)} warning | "
        f"(overdue ACR: {_a(totals.get('overdueAcr', 0))}/mo, "
        f"critical ACR: {_a(totals.get('criticalAcr', 0))}/mo)",
        "",
        "### Overdue & Critical Accounts",
        "| Account | Directors | Status | Days | ACR/mo | At-Risk Models |",
        "|---------|-----------|--------|------|--------|----------------|",
    ]
    for a in urgent:
        emoji = "⚫" if a["level"] == "overdue" else "🔴"
        days_val = a.get("minDays")
        days_str = "OVERDUE" if days_val is None or days_val <= 0 else str(days_val)
        models = ", ".join(a.get("atRiskModels", [])[:3])
        dirs = ", ".join(a.get("directors", []))
        ctx.append(
            f"| {a['name']} | {dirs} | {emoji} {a['level'].upper()} | "
            f"{days_str} | {_a(a.get('monthlyAcr', 0))} | {models} |"
        )

    if warning_accts:
        ctx += [
            "",
            "### Warning Accounts (≤180 days)",
            "| Account | Directors | Days | ACR/mo | At-Risk Models |",
            "|---------|-----------|------|--------|----------------|",
        ]
        for a in warning_accts:
            models = ", ".join(a.get("atRiskModels", [])[:2])
            dirs = ", ".join(a.get("directors", []))
            ctx.append(
                f"| {a['name']} | {dirs} | {a.get('minDays', '?')} | "
                f"{_a(a.get('monthlyAcr', 0))} | {models} |"
            )

    if model_summary:
        ctx += [
            "",
            "### Fleet-Wide Model Retirements (≤180 days)",
            "| Model | Ret. Date | Days | Accts | Deploys | Replacement |",
            "|-------|-----------|------|-------|---------|-------------|",
        ]
        for m in model_summary[:20]:
            ctx.append(
                f"| {m['model']} | {m['retirementDate']} | {m['days']} | "
                f"{m['hlsAccounts']} | {m['hlsDeploys']} | {m['replacement']} |"
            )

    ctx += [
        "",
        "### Director Summary",
        "| Director | Accounts | ACR/mo | ⚫ Overdue | 🔴 Critical | 🟡 Warning |",
        "|----------|----------|--------|-----------|------------|-----------|",
    ]
    for d in director_summaries[:15]:
        ctx.append(
            f"| {d['director']} | {d['accounts']} | {_a(d['monthlyAcr'])} | "
            f"{d['overdue']} | {d['critical']} | {d['warning']} |"
        )

    context = "\n".join(ctx)

    system = (
        "You are a strategic advisor for Microsoft's HLS (Health & Life Sciences) "
        "Customer Success organization. Generate specific, actionable recommendations "
        "based on Azure OpenAI model retirement risk data. Name exact accounts, models, "
        "and CSA managers. Prioritize by revenue (ACR) and urgency. "
        "Output clean GitHub-flavored markdown. No preamble or meta-commentary."
    )

    user = f"""Based on the HLS CSA org tracker data below, generate a recommendations report.

{context}

Produce a markdown report with these exact sections:

# HLS CSA Model IQ — Recommendations
*Generated {today.isoformat()}*

## 🎯 Executive Recommendations
3-5 bullet-point priority actions for org leadership. Lead each with the account or model name and ACR at risk.

## ⚡ Immediate Actions Required
For every OVERDUE and CRITICAL account (list all of them): specific migration action, responsible director/CSA alias, model to migrate FROM → TO, and a realistic timeline. Group by director. Include ACR at stake.

## 📋 Priority Migration Plans
Top 10 at-risk accounts by ACR: for each, one paragraph with current at-risk models, recommended replacement model(s), migration complexity (Low/Medium/High), and suggested completion date.

## 🔄 Model Fleet Migration Roadmap
For each retiring model in the fleet: recommended replacement (use known Azure OpenAI successors: gpt-4o → gpt-4.1, gpt-4o-mini → gpt-4.1-mini, gpt-35-turbo → gpt-4o-mini), migration notes (prompt compatibility, latency/cost tradeoffs), affected HLS account count, and retirement deadline.

## 📊 Director Action Plans
For each director with at-risk accounts: bullet list of their accounts requiring action with specific next steps and timeline.

Use exact account names, model names, and director aliases from the data. Be specific and actionable."""

    client = get_client()
    deployment = get_deployment("architecture")

    resp = client.chat.completions.create(
        model=deployment,
        messages=[
            {"role": "system", "content": system},
            {"role": "user", "content": user},
        ],
        temperature=0.3,
        max_completion_tokens=4000,
    )
    return resp.choices[0].message.content or ""


# ── PDF export ────────────────────────────────────────────────────────────────

_PDF_CSS = """
<style>
  body {
    font-family: 'Segoe UI', -apple-system, BlinkMacSystemFont, sans-serif;
    color: #1a1a1a;
    background: #ffffff;
    max-width: 960px;
    margin: 0 auto;
    padding: 32px 24px;
    line-height: 1.6;
    font-size: 14px;
  }
  h1 { color: #0078d4; font-size: 24px; border-bottom: 2px solid #0078d4; padding-bottom: 8px; }
  h2 { color: #106ebe; font-size: 20px; margin-top: 28px; border-bottom: 1px solid #e0e0e0; padding-bottom: 6px; }
  h3 { color: #323130; font-size: 16px; margin-top: 20px; }
  h4 { color: #605e5c; font-size: 14px; margin-top: 16px; }
  table { border-collapse: collapse; width: 100%; margin: 16px 0; font-size: 13px; }
  th {
    background-color: #0078d4; color: #ffffff; font-weight: 600;
    text-align: left; padding: 10px 12px; white-space: nowrap;
    border: 1px solid #005a9e;
  }
  td { padding: 8px 12px; border: 1px solid #e0e0e0; vertical-align: top; }
  tr:nth-child(even) td { background-color: #f8f9fa; }
  td:first-child { font-weight: 500; }
  code {
    background-color: #f4f4f4; padding: 2px 6px; border-radius: 3px;
    font-family: 'Cascadia Code', 'Consolas', monospace; font-size: 12px;
  }
  pre { background-color: #1e1e1e; color: #d4d4d4; padding: 16px; border-radius: 6px; overflow-x: auto; font-size: 12px; }
  pre code { background: none; padding: 0; color: inherit; }
  blockquote { border-left: 4px solid #0078d4; margin: 16px 0; padding: 8px 16px; background-color: #f0f6ff; color: #323130; }
  a { color: #0078d4; text-decoration: none; }
  ul, ol { padding-left: 24px; }
  li { margin: 4px 0; }
  hr { border: none; border-top: 1px solid #e0e0e0; margin: 24px 0; }
  strong { color: #323130; }
  .report-footer { margin-top: 40px; padding-top: 16px; border-top: 1px solid #e0e0e0; font-size: 12px; color: #605e5c; }
  @media print {
    body { max-width: 100%; padding: 0; }
    table { page-break-inside: avoid; }
    h2 { page-break-after: avoid; }
  }
</style>
"""

_PDF_FOOTER = """
<div class="report-footer">
  Generated by <strong>Azure Architect AI</strong> — HLS CSA Org Tracker<br>
  <em>This report was auto-generated from uploaded data files.</em>
</div>
"""


def build_org_report_pdf(md_text: str, generated: str = "") -> bytes:
    """Convert the org tracker markdown to PDF via Playwright Chromium."""
    import tempfile

    import markdown as md_lib
    from playwright.sync_api import sync_playwright

    html_body = md_lib.markdown(
        md_text,
        extensions=["tables", "fenced_code", "toc", "nl2br"],
    )
    title = "HLS CSA Org Tracker"
    html_doc = f"""<!DOCTYPE html>
<html lang="en">
<head>
  <meta charset="UTF-8">
  <meta name="viewport" content="width=device-width, initial-scale=1.0">
  <title>{title}</title>
  {_PDF_CSS}
</head>
<body>
{html_body}
{_PDF_FOOTER}
</body>
</html>"""

    with tempfile.NamedTemporaryFile(suffix=".html", delete=False, mode="w", encoding="utf-8") as tmp:
        tmp.write(html_doc)
        tmp_path = Path(tmp.name)

    try:
        with sync_playwright() as p:
            browser = p.chromium.launch()
            page = browser.new_page()
            page.goto(tmp_path.absolute().as_uri(), wait_until="networkidle")
            pdf_bytes = page.pdf(
                format="A4",
                margin={"top": "20mm", "bottom": "20mm", "left": "15mm", "right": "15mm"},
                print_background=True,
            )
            browser.close()
    finally:
        tmp_path.unlink(missing_ok=True)

    return pdf_bytes


# ── PPTX color palette ─────────────────────────────────────────────────────
_C_DARK   = (0x35, 0x48, 0x5E)   # soft blue-grey header
_C_WHITE  = (0xFF, 0xFF, 0xFF)
_C_LTBLUE = (0x5B, 0xA3, 0xD0)   # muted sky-blue accent strip
_C_MSFT   = (0x2B, 0x7A, 0xC8)   # softer Microsoft blue
_C_RED    = (0xBE, 0x45, 0x45)   # muted rose-red (overdue)
_C_ORANGE = (0xCF, 0x7A, 0x3A)   # muted amber-orange (critical)
_C_AMBER  = (0xC9, 0xA0, 0x2E)   # muted gold (warning)
_C_GRAY   = (0x7A, 0x8D, 0x9C)   # muted blue-grey text
_C_SLATE  = (0x4F, 0x62, 0x74)   # medium slate (table headers)
_C_LTBG   = (0xF7, 0xF8, 0xFA)   # near-white slide background

_RISK_COLORS: dict[str, tuple[int, int, int]] = {
    "overdue":  _C_RED,
    "critical": _C_ORANGE,
    "warning":  _C_AMBER,
    "ok":       _C_GRAY,
}


def make_org_data(org_scorecard: dict[str, Any], model_summary: list[dict[str, Any]], month_col: str) -> dict[str, Any]:
    """Return a JSON-safe dict suitable for sending to the frontend and back."""
    def _default(obj: object) -> str:
        if isinstance(obj, date):
            return obj.isoformat()
        raise TypeError(f"Not serializable: {type(obj)}")

    raw: dict[str, Any] = {
        "totals": org_scorecard.get("totals", {}),
        "allAccounts": org_scorecard.get("allAccounts", []),
        "model_summary": model_summary,
        "month_label": month_col,
    }
    return json.loads(json.dumps(raw, default=_default))  # type: ignore[no-any-return]


def build_org_report_pptx(org_data: dict[str, Any], today: date) -> bytes:
    """Generate a 4-slide Model IQ Risk Review PPTX matching the reference deck style."""
    from io import BytesIO

    from pptx import Presentation  # type: ignore[import-untyped]
    from pptx.dml.color import RGBColor  # type: ignore[import-untyped]
    from pptx.util import Emu, Pt  # type: ignore[import-untyped]

    EMU_PER_IN = 914400

    def _rgb(t: tuple[int, int, int]) -> RGBColor:
        return RGBColor(*t)

    def _inches(v: float) -> Emu:
        return Emu(int(v * EMU_PER_IN))

    def _rect(slide: Any, left: float, t: float, w: float, h: float,
              fill: tuple[int, int, int] | None = None, alpha: int | None = None) -> Any:
        from pptx.util import Emu as E2
        shape = slide.shapes.add_shape(
            1,  # MSO_AUTO_SHAPE_TYPE.RECTANGLE
            E2(int(left * EMU_PER_IN)), E2(int(t * EMU_PER_IN)),
            E2(int(w * EMU_PER_IN)), E2(int(h * EMU_PER_IN)),
        )
        shape.line.fill.background()
        if fill:
            shape.fill.solid()
            shape.fill.fore_color.rgb = _rgb(fill)
        else:
            shape.fill.background()
        return shape

    def _add_text(slide: Any, text: str, left: float, t: float, w: float, h: float,
                  size: float = 12, bold: bool = False, color: tuple[int, int, int] = _C_WHITE,
                  align: str = "left", wrap: bool = True) -> Any:
        from pptx.enum.text import PP_ALIGN  # type: ignore[import-untyped]
        from pptx.util import Emu as E2
        tb = slide.shapes.add_textbox(
            E2(int(left * EMU_PER_IN)), E2(int(t * EMU_PER_IN)),
            E2(int(w * EMU_PER_IN)), E2(int(h * EMU_PER_IN)),
        )
        tb.text_frame.word_wrap = wrap
        p = tb.text_frame.paragraphs[0]
        al = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}.get(align, PP_ALIGN.LEFT)
        p.alignment = al
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = _rgb(color)
        return tb

    def _header(slide: Any, title: str, subtitle: str = "") -> None:
        _rect(slide, 0, 0, 10, 0.50, fill=_C_DARK)
        _rect(slide, 0, 0.50, 10, 0.05, fill=_C_LTBLUE)
        _add_text(slide, title, 0.2, 0.05, 6.0, 0.45, size=18, bold=True, color=_C_WHITE)
        if subtitle:
            _add_text(slide, subtitle, 6.1, 0.08, 3.7, 0.38, size=10, bold=False, color=_C_LTBLUE, align="right")
        _add_text(
            slide,
            f"Azure HLS CSA  |  Model IQ Risk Review  |  {today.strftime('%B %Y')}",
            0.1, 5.40, 9.8, 0.2, size=7, bold=False, color=_C_GRAY, align="center",
        )

    def _table_header_row(table: Any, headers: list[str],
                          fill: tuple[int, int, int] = _C_SLATE,
                          txt: tuple[int, int, int] = _C_WHITE, size: float = 8) -> None:
        from pptx.enum.text import PP_ALIGN  # type: ignore[import-untyped]
        row = table.rows[0]
        for i, hdr in enumerate(headers):
            cell = row.cells[i]
            cell.fill.solid()
            cell.fill.fore_color.rgb = _rgb(fill)
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = hdr
            run.font.size = Pt(size)
            run.font.bold = True
            run.font.color.rgb = _rgb(txt)

    def _table_data_row(table: Any, row_idx: int, values: list[str],
                        fill: tuple[int, int, int] | None = None,
                        txt: tuple[int, int, int] = _C_DARK, size: float = 7.5,
                        aligns: list[str] | None = None) -> None:
        from pptx.enum.text import PP_ALIGN  # type: ignore[import-untyped]
        row = table.rows[row_idx]
        for i, val in enumerate(values):
            cell = row.cells[i]
            if fill:
                cell.fill.solid()
                cell.fill.fore_color.rgb = _rgb(fill)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = _rgb(_C_WHITE if row_idx % 2 == 0 else _C_LTBG)
            p = cell.text_frame.paragraphs[0]
            al_str = (aligns[i] if aligns and i < len(aligns) else "left")
            p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}.get(al_str, PP_ALIGN.LEFT)
            run = p.add_run()
            run.text = val
            run.font.size = Pt(size)
            run.font.color.rgb = _rgb(txt)

    # ── Presentation setup ────────────────────────────────────────────────
    prs = Presentation()
    prs.slide_width  = _inches(10.0)
    prs.slide_height = _inches(5.63)
    blank_layout = prs.slide_layouts[6]

    totals       = org_data.get("totals", {})
    all_accounts = org_data.get("allAccounts", [])
    model_sum    = org_data.get("model_summary", [])
    month_label  = org_data.get("month_label", "")

    overdue_accounts  = [a for a in all_accounts if a.get("level") == "overdue"]
    critical_accounts = [a for a in all_accounts if a.get("level") == "critical"]

    def _fmt_acr(v: float) -> str:
        if v >= 1_000_000:
            return f"${v / 1_000_000:.1f}M"
        if v >= 1_000:
            return f"${v / 1_000:.0f}K"
        return f"${v:.0f}"

    def _fmt_days(d: int | None) -> str:
        if d is None:
            return "—"
        if d <= 0:
            return f"{abs(d)}d overdue"
        return f"{d}d"

    # ── Slide 1: Portfolio at a Glance ────────────────────────────────────
    s1 = prs.slides.add_slide(blank_layout)
    _rect(s1, 0, 0, 10, 5.63, fill=_C_LTBG)
    _header(s1, "Portfolio at a Glance", subtitle=f"ACR: {month_label}")

    total_acr  = totals.get("totalMonthlyAcr", 0)
    n_accounts = totals.get("accountsWithDeployments", 0)
    n_deploys  = totals.get("totalDeployments", 0)
    n_overdue  = totals.get("overdue", 0)
    n_critical = totals.get("critical", 0)
    n_warning  = totals.get("warning", 0)

    kpi_cards = [
        ("Total Monthly ACR", _fmt_acr(total_acr)),
        ("Accounts", str(n_accounts)),
        ("AI Deployments", str(n_deploys)),
    ]
    card_w, card_h = 2.5, 1.0
    for i, (label, val) in enumerate(kpi_cards):
        cx = 0.5 + i * 3.1
        _rect(s1, cx, 0.65, card_w, card_h, fill=_C_DARK)
        _add_text(s1, label, cx + 0.1, 0.67, card_w - 0.2, 0.3, size=7, bold=False, color=_C_LTBLUE, align="center")
        _add_text(s1, val,   cx + 0.1, 0.97, card_w - 0.2, 0.55, size=20, bold=True, color=_C_WHITE, align="center")

    risk_panels = [
        ("OVERDUE", n_overdue, totals.get("overdueAcr", 0),  _C_RED,    2.6),
        ("CRITICAL", n_critical, totals.get("criticalAcr", 0), _C_ORANGE, 5.5),
    ]
    for label, count, acr, col, cx in risk_panels:
        _rect(s1, cx, 1.80, 2.35, 1.20, fill=col)
        _add_text(s1, label,         cx + 0.1, 1.82, 2.15, 0.28, size=9, bold=True, color=_C_WHITE, align="center")
        _add_text(s1, str(count),    cx + 0.1, 2.10, 2.15, 0.55, size=22, bold=True, color=_C_WHITE, align="center")
        _add_text(s1, "accounts",    cx + 0.1, 2.62, 2.15, 0.22, size=7, bold=False, color=_C_WHITE, align="center")
        _add_text(s1, _fmt_acr(acr) + " at risk", cx + 0.1, 2.82, 2.15, 0.22, size=7, bold=False, color=_C_WHITE, align="center")

    _rect(s1, 0.4, 1.78, 2.0, 1.24, fill=_C_MSFT)
    _add_text(s1, "WARNING",      0.5, 1.80, 1.8, 0.28, size=9, bold=True, color=_C_WHITE, align="center")
    _add_text(s1, str(n_warning), 0.5, 2.08, 1.8, 0.55, size=22, bold=True, color=_C_WHITE, align="center")
    _add_text(s1, "accounts",     0.5, 2.60, 1.8, 0.22, size=7, bold=False, color=_C_WHITE, align="center")

    # Top 5 accounts table
    top5 = sorted(all_accounts, key=lambda a: a.get("monthlyAcr", 0), reverse=True)[:5]
    if top5:
        top5_headers = ["Account", "ACR/mo", "Risk", "Models at Risk"]
        col_w_top5 = [3.0, 0.9, 0.8, 2.8]
        tbl_top5 = s1.shapes.add_table(len(top5) + 1, 4,
                                        _inches(0.2), _inches(3.15),
                                        _inches(9.6), _inches(2.1)).table
        for ci, cw in enumerate(col_w_top5):
            tbl_top5.columns[ci].width = _inches(cw)
        _table_header_row(tbl_top5, top5_headers, fill=_C_SLATE)
        for ri, acct in enumerate(top5):
            risk_lvl = acct.get("level", "ok")
            rc = _RISK_COLORS.get(risk_lvl, _C_GRAY)
            models = ", ".join((acct.get("atRiskModels") or [])[:3])
            _table_data_row(tbl_top5, ri + 1, [
                acct.get("name", ""),
                _fmt_acr(acct.get("monthlyAcr", 0)),
                risk_lvl.upper(),
                models or "—",
            ], aligns=["left", "right", "center", "left"])
            days = acct.get("minDays")
            if days is not None and days <= 0:
                tbl_top5.rows[ri + 1].cells[2].fill.solid()
                tbl_top5.rows[ri + 1].cells[2].fill.fore_color.rgb = _rgb(rc)

    # ── Slide 2: Retiring Models ───────────────────────────────────────────
    s2 = prs.slides.add_slide(blank_layout)
    _rect(s2, 0, 0, 10, 5.63, fill=_C_LTBG)
    _header(s2, "Models Retiring Within 180 Days", subtitle="Fleet-Wide Exposure")

    retiring = [m for m in model_sum if (m.get("days") or 999) <= 180]
    retiring.sort(key=lambda m: m.get("days") or 999)
    headers2 = ["Model", "Retirement Date", "Days", "HLS Deployments", "HLS Accounts", "Replacement"]
    col_w2   = [2.0, 1.4, 0.65, 1.35, 1.2, 3.0]
    n_rows2  = max(len(retiring), 1)
    tbl2 = s2.shapes.add_table(n_rows2 + 1, 6,
                                _inches(0.2), _inches(0.65),
                                _inches(9.6), _inches(4.7)).table
    for ci, cw in enumerate(col_w2):
        tbl2.columns[ci].width = _inches(cw)
    _table_header_row(tbl2, headers2, fill=_C_DARK)
    if not retiring:
        _table_data_row(tbl2, 1, ["No models retiring within 180 days", "", "", "", "", ""])
    else:
        for ri, m in enumerate(retiring[:n_rows2]):
            days = m.get("days")
            if days is not None and days <= 0:
                lvl_col = _C_RED
            elif days is not None and days <= 90:
                lvl_col = _C_ORANGE
            elif days is not None and days <= 180:
                lvl_col = _C_AMBER
            else:
                lvl_col = None
            _table_data_row(tbl2, ri + 1, [
                m.get("model", ""),
                m.get("retirementDate", ""),
                _fmt_days(days),
                str(m.get("hlsDeploys", 0)),
                str(m.get("hlsAccounts", 0)),
                m.get("replacement", "—") or "—",
            ], fill=lvl_col, txt=_C_WHITE if lvl_col else _C_DARK,
               aligns=["left", "center", "center", "center", "center", "left"])

    # ── Slide 3: Overdue Accounts ──────────────────────────────────────────
    s3 = prs.slides.add_slide(blank_layout)
    _rect(s3, 0, 0, 10, 5.63, fill=_C_LTBG)
    _header(s3, "Overdue Accounts", subtitle="Retirement date has passed")

    headers3 = ["Account", "Director", "Monthly ACR", "Days Overdue", "Models at Risk"]
    col_w3   = [2.8, 1.5, 1.1, 1.1, 4.1]
    n_rows3  = max(len(overdue_accounts), 1)
    tbl3 = s3.shapes.add_table(n_rows3 + 1, 5,
                                _inches(0.2), _inches(0.65),
                                _inches(9.6), _inches(4.7)).table
    for ci, cw in enumerate(col_w3):
        tbl3.columns[ci].width = _inches(cw)
    _table_header_row(tbl3, headers3, fill=_C_RED)
    if not overdue_accounts:
        _table_data_row(tbl3, 1, ["No overdue accounts", "", "", "", ""])
    else:
        for ri, acct in enumerate(overdue_accounts[:n_rows3]):
            dirs = ", ".join(acct.get("directors") or [])
            models = ", ".join((acct.get("atRiskModels") or [])[:4])
            days = acct.get("minDays")
            _table_data_row(tbl3, ri + 1, [
                acct.get("name", ""),
                dirs,
                _fmt_acr(acct.get("monthlyAcr", 0)),
                _fmt_days(days),
                models or "—",
            ], aligns=["left", "left", "right", "center", "left"])

    # ── Slide 4: Action Plan ───────────────────────────────────────────────
    s4 = prs.slides.add_slide(blank_layout)
    _rect(s4, 0, 0, 10, 5.63, fill=_C_LTBG)
    _header(s4, "Action Plan", subtitle="Immediate and Near-Term Priorities")

    # Left column — overdue
    _rect(s4, 0.2, 0.65, 4.65, 0.35, fill=_C_RED)
    _add_text(s4, f"OVERDUE — {len(overdue_accounts)} Accounts", 0.3, 0.67, 4.45, 0.30,
              size=10, bold=True, color=_C_WHITE)
    col_left_headers = ["Account", "ACR/mo", "Key Models"]
    col_left_w = [2.2, 0.9, 1.55]
    n_od = max(len(overdue_accounts), 1)
    tbl4l = s4.shapes.add_table(n_od + 1, 3,
                                 _inches(0.2), _inches(1.0),
                                 _inches(4.65), _inches(4.3)).table
    for ci, cw in enumerate(col_left_w):
        tbl4l.columns[ci].width = _inches(cw)
    _table_header_row(tbl4l, col_left_headers, fill=_C_SLATE, size=7.5)
    if not overdue_accounts:
        _table_data_row(tbl4l, 1, ["No overdue accounts", "", ""])
    else:
        for ri, acct in enumerate(overdue_accounts[:n_od]):
            models = ", ".join((acct.get("atRiskModels") or [])[:2])
            _table_data_row(tbl4l, ri + 1, [
                acct.get("name", ""),
                _fmt_acr(acct.get("monthlyAcr", 0)),
                models or "—",
            ], aligns=["left", "right", "left"])

    # Right column — critical
    _rect(s4, 5.15, 0.65, 4.65, 0.35, fill=_C_ORANGE)
    _add_text(s4, f"CRITICAL — {len(critical_accounts)} Accounts (≤90 days)", 5.25, 0.67, 4.45, 0.30,
              size=10, bold=True, color=_C_WHITE)
    n_cr = max(len(critical_accounts), 1)
    tbl4r = s4.shapes.add_table(n_cr + 1, 3,
                                 _inches(5.15), _inches(1.0),
                                 _inches(4.65), _inches(4.3)).table
    for ci, cw in enumerate(col_left_w):
        tbl4r.columns[ci].width = _inches(cw)
    _table_header_row(tbl4r, col_left_headers, fill=_C_SLATE, size=7.5)
    if not critical_accounts:
        _table_data_row(tbl4r, 1, ["No critical accounts", "", ""])
    else:
        for ri, acct in enumerate(critical_accounts[:n_cr]):
            models = ", ".join((acct.get("atRiskModels") or [])[:2])
            _table_data_row(tbl4r, ri + 1, [
                acct.get("name", ""),
                _fmt_acr(acct.get("monthlyAcr", 0)),
                models or "—",
            ], aligns=["left", "right", "left"])

    out = BytesIO()
    prs.save(out)
    return out.getvalue()


def build_download_all_zip(
    markdown: str,
    recommendations_markdown: str,
    generated: str,
    org_data: dict[str, Any],
    today: date,
) -> bytes:
    """Bundle report PDF, recommendations PDF (if any), and PPTX into a single ZIP."""
    import zipfile

    date_str = today.isoformat()
    report_pdf = build_org_report_pdf(markdown, generated)
    pptx_bytes = build_org_report_pptx(org_data, today)

    out = io.BytesIO()
    with zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED) as zf:
        zf.writestr(f"hls-csa-org-tracker-{date_str}.pdf", report_pdf)
        if recommendations_markdown.strip():
            recs_pdf = build_org_report_pdf(recommendations_markdown, generated)
            zf.writestr(f"hls-csa-recommendations-{date_str}.pdf", recs_pdf)
        zf.writestr(f"hls-csa-model-iq-{date_str}.pptx", pptx_bytes)
    return out.getvalue()
